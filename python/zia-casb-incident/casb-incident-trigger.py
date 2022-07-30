#!/usr/bin/python
'''

This script is to generate documents containing data that can trigger DLP policies, and automatically put them in OneDrive or Box.
It will also copy a random chosen malware sample to oneDrive or Box with a random name. The malware does already have to be present on the machine this is run on

The DLP files that are generated are randomized with the following parameters:
    - Dataset: The actual data in these sets are in lists variables. Currently the datasets in here are:
        - CCN
        - SSN
        - PCI
        - Safemarch Confidential

    - File Format: Currently the script supports creating:
        - pdf
        - txt
        - csv
        - docx
        - xlsx
        - pptx

    - Directory in OneDrive or Box. These already have to exist and are the following:
        - Personal
        - Planning
        - Presentations
        - Projects
        - QBR
        - Roadmap

    - A file name chosen based on the dataset (extension is chosen by fileformat)
        - CCN: CompanyCreditCards, CustomerData, PaymentDetails, Transactions, CorpAccounts
        - SSN: EmployeeData, HR, CurrentEmployees, CustomerInfo, VerificationList
        - PCI: SalesForceData, PointOfSaleData, Backup, Export, CustomerDetails
        - Safemarch Confidential: Roadmap, QBR, Q1Results, SurveyResults, ProjectZenith

From these variables, files get created on the fly using the data from the datasets. It will also pad each file with 256 random characters, so every file is different.

For best results (multiple users), run this on a machine with OneDrive and Box clients installed. OneDrive will let you add 10 drives at once.
Though one drive is reserved for OneDrive Personal... Box will let you add one. So you can have 10 business users.
For triggering and doing multiple actions, it's possible to configure OneDrive to automatically link share publicly a file when it appears. Look for "Power Automate"

It also checks if it is during working hours or weekend and will throttle generation.


This script is copyleft. NO WARRANTIES EXPRESSED OR IMPLIED.
'''


import docx, random, openpyxl, string, time, os, shutil
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime

# Drive locations - Directories that map to OneDrive or Box. Put in a list so we can randomly choose one. Always include trailing \\
# Note that when choosing a folder for OneDrive, OneDrive automatically adds "oneDrive - CompanyName" as a subdir
driveLocation = ([''
                'c:\\onedrive\\oneDrive - Contoso\\',
                'c:\\onedrive1\\oneDrive - Contoso\\',
                'c:\\onedrive2\\oneDrive - Contoso\\',
                'c:\\onedrive3\\oneDrive - Contoso\\',
                'c:\\onedrive4\\oneDrive - Contoso\\',
                'c:\\onedrive5\\oneDrive - Contoso\\',
                'c:\\onedrive6\\oneDrive - Contoso\\',
                'c:\\onedrive7\\oneDrive - Contoso\\',
                'c:\\onedrive8\\oneDrive - Contoso\\',
                'c:\\onedrive9\\oneDrive - Contoso\\',
                'c:\\users\z-tm\\box\\'
                ])


# This is the location where the script will pull files to use as randomized malware. Have at least a single file in here
malwareLocation = 'c:\\tfgen\\malware\\'



# datasets
CCN = ['CCN', '4716653795988594', '4556147740704506', '4820497634813420', '2223000010476528', '2223000010476510',
       '2223000048400011', '6011792346848217', '6011543730947358', '6011038898889385']
SSN = ['SSN', 'SSN Number - 234-34-5026', 'SSN Number - 242-42-4620', 'SSN Number - 345-22-5678',
       'SSN Number - 548-98-9743', 'SSN Number - 205-70-0003', 'SSN Number - 436-88-7865', 'SSN Number - 477-44-4494']
PCI = ['PCI', 'First and Last Name,SSN,Credit Card Number (VISA, MC, AMEX)',
       'Robert Aragon,489-36-8350,4929-3813-3266-4295', 'Ashley Borden,514-14-8905,5370-4638-8881-3020',
       'Thomas Conley,690-05-5315,4916-4811-5814-8111', 'Susan Davis,421-37-1396,4916-4034-9269-8783',
       'Christopher Diaz,458-02-6124,5299-1561-5689-1938', 'Rick Edwards,612-20-6832,5293-8502-0071-3058',
       'Victor Faulkner,300-62-3266,5548-0246-6336-5664', 'Lisa Garrison,660-03-8360,4539-5385-7425-5825',
       'Marjorie Green,213-46-8915,4916-9766-5240-6147', 'Mark Hall,449-48-3135,4556-0072-1294-7415',
       'James Heard,559-81-1301,4532-4220-6922-9909', 'Albert Iorio,322-84-2281,4916-6734-7572-5015',
       'Charles Jackson,646-44-9061,5218-0144-2703-9266', 'Teresa Kaminski,465-73-5022,5399-0706-4128-0178',
       'Tim Lowe,044-34-6954,5144-8691-2776-1108']
Confidential = ['Confidential', 'Safemarch Confidential', 'Do not distribute', 'Internal Only']


# This function creates a random string, with the parameter being the number of characters
def randomString(length):
    letters = string.ascii_letters
    return ''.join(random.choice(letters) for i in range(length))


# This function is to generate a docx file, with as input dataset and directory. It will also call the determineFileName function to get the filename
# If the filename in that directory already exists, delete it first. New file will have same name, but different random data
def createDoc(dataSet, directory):
    fileName = determineFileName(dataset)
    if os.path.isfile(directory + fileName + '.docx'):
        os.unlink(directory + fileName + '.docx')
    newFile = docx.Document()
    for i in range(1, len(dataset)):
        newFile.add_paragraph(dataset[i])
    newFile.add_paragraph(randomString(256))
    newFile.save(directory + fileName + '.docx')
    print('creating docx with filename: ' + fileName + '.docx in directory: ' + directory + ' containing ' + dataset[0] + ' data\n')
    return


# This function is to generate a csv file, with as input dataset and directory. It will also call the determineFileName function to get the filename
# If the filename in that directory already exists, delete it first. New file will have same name, but different random data
def createCsv(dataset, directory):
    fileName = determineFileName(dataset)
    if os.path.isfile(directory + fileName + '.csv'):
        os.unlink(directory + fileName + '.csv')
    newFile = open(directory + fileName + '.csv', 'w')
    for i in range(1, len(dataset)):
        newFile.write(dataset[i] + '\n')
    newFile.write(randomString(256))
    newFile.close()
    print('creating csv with filename: ' + fileName + '.csv in directory: ' + directory + ' containing ' + dataset[
        0] + ' data\n')
    return



# This function is to generate a txt file, with as input dataset and directory. It will also call the determineFileName function to get the filename
# If the filename in that directory already exists, delete it first. New file will have same name, but different random data
def createTxt(dataset, directory):
    fileName = determineFileName(dataset)
    if os.path.isfile(directory + fileName + '.txt'):
        os.unlink(directory + fileName + '.txt')
    newFile = open(directory + fileName + '.txt', 'w')
    for i in range(1, len(dataset)):
        newFile.write(dataset[i] + '\n')
    newFile.write(randomString(256))
    newFile.close()
    print('creating txt with filename: ' + fileName + '.txt in directory: ' + directory + ' containing ' + dataset[
        0] + ' data\n')
    return


# This function is to generate a xlsx file, with as input dataset and directory. It will also call the determineFileName function to get the filename
# If the filename in that directory already exists, delete it first. New file will have same name, but different random data
def createXls(dataset, directory):
    fileName = determineFileName(dataset)
    if os.path.isfile(directory + fileName + '.xlsx'):
        os.unlink(directory + fileName + '.xlsx')
    newFile = openpyxl.Workbook()
    newFile.sheetnames
    sheet = newFile.active
    sheet.title = 'Confidential'
    for i in range(1, len(dataset)):
        sheet['A' + str(i)] = dataset[i]
    sheet['A100'] = randomString(256)
    newFile.save(directory + fileName + '.xlsx')
    print('creating xlsx with filename: ' + fileName + '.xlsx in directory: ' + directory + ' containing ' + dataset[0] + ' data\n')
    return

# This function is to generate a pdf file, with as input dataset and directory. It will also call the determineFileName function to get the filename
# If the filename in that directory already exists, delete it first. New file will have same name, but different random data
def createPdf(dataset, directory):
    fileName = determineFileName(dataset)
    if os.path.isfile(directory + fileName + '.pdf'):
        os.unlink(directory + fileName + '.pdf')
    pdf = FPDF()
    pdf.set_font('Arial', size=12)
    pdf.add_page()
    for i in range(1,len(dataset)):
        pdf.cell(2000,25, txt=dataset[i], ln=1)
    pdf.cell(2000,25, txt=randomString(256))
    pdf.output(directory + fileName + '.pdf')
    print('creating pdf with filename: ' + fileName + '.pdf in directory: ' + directory + ' containing ' + dataset[0] + ' data\n')
    return


# This function is to generate a pptx file, with as input dataset and directory. It will also call the determineFileName function to get the filename
# If the filename in that directory already exists, delete it first. New file will have same name, but different random data
def createPpt(dataset, directory):
    fileName = determineFileName(dataset)
    if os.path.isfile(directory + fileName + '.pptx'):
        os.unlink(directory + fileName + '.pptx')
    ppt = Presentation()
    blank_slide_layout = ppt.slide_layouts[6]
    slide = ppt.slides.add_slide(blank_slide_layout)
    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    for i in range(1,len(dataset)):

        p = tf.add_paragraph()
        p.text = dataset[i]
    p =tf.add_paragraph()
    p.text = randomString(256)
    ppt.save(directory + fileName + '.pptx')
    print('creating pptx with filename: ' + fileName + '.pptx in directory: ' + directory + ' containing ' + dataset[
        0] + ' data\n')
    return

# This function returns a filename to use based on the dataset input. The filename get randomly chosen based on the dataset
def determineFileName(dataset):
    if dataset[0] == 'CCN':
        #print('data is CCN')
        possibleNames = ['CompanyCreditCards', 'CustomerData', 'PaymentDetails', 'Transactions', 'CorpAccounts']
        fileName = random.choice(possibleNames)
    elif dataset[0] == 'SSN':
        #print('data is SSN')
        possibleNames = ['EmployeeData', 'HR', 'CurrentEmployees', 'CustomerInfo', 'VerificationList']
        fileName = random.choice(possibleNames)
    elif dataset[0] == 'PCI':
        #print('data is PCI')
        possibleNames = ['SalesForceData', 'PointOfSaleData', 'Backup', 'Export', 'CustomerDetails']
        fileName = random.choice(possibleNames)
    elif dataset[0] == 'Confidential':
        #print('data is Confidential')
        possibleNames = ['Roadmap', 'QBR', 'Q1Results', 'SurveyResults', 'ProjectZenith']
        fileName = random.choice(possibleNames)

    return fileName

# This is the function that will determine which dataset to use, as well as which oneDrive/Box user to use, as well as which subdir in the user's OneDrive/Box Account and the fileformat
# It then returns the fullpath, as well as the dataset to use and the fileformat
def pickRandomDataAndFile():
    data = [CCN, SSN, PCI, Confidential]
    fileFormat = [createDoc, createCsv, createTxt, createXls, createPdf, createPpt]
    Directories = ['Personal', 'Presentations', 'Projects', 'Planning', 'QBR', 'Roadmap']
    dataset = random.choice(data)
    file = random.choice(fileFormat)
    directory = random.choice(Directories)
    fullPath = driveLocation[random.randint(0,len(driveLocation)-1)] + directory + '\\'

    return dataset, file, fullPath

# This function will copy a random file from malwareLocation to a random directory with a random name
def putRandomMalware():
    malware = random.choice(os.listdir(malwareLocation))
    dataset, file, directory = pickRandomDataAndFile()
    malwareFileName = randomString(8) + '.bin'
    shutil.copy(malwareLocation + malware, directory + malwareFileName)
    print('Copying malware file: '+ malwareFileName + ' to: ' + directory + malwareFileName + '\n')
    return


# This script will run indefinitely until canceled (or breaks). It also includes the timers to keep generation at manageable pace
# Without the sleeps, files get created at a breakneck speed. might be good for initial seeding. It uses hour of the day and day of the week to determine the sleep times
# Sleep timers are random inside ranges, to make it look more natural.
# To increase or decrease malware occurrence, chance "0.05" to whatever you want. 0.05 = 5% of the time
while True:
    dataset, file, directory = pickRandomDataAndFile()

    # call the function that was chosen with pickRandomDataAndFile and pass dataset and directory
    file(dataset, directory)

    if random.random() < 0.05:
        putRandomMalware()


    if (18 <= datetime.now().hour <= 8) or (datetime.today().weekday() > 5):
        time.sleep(random.randint(7200, 25000))
    else:
        time.sleep(random.randint(1800,3600))